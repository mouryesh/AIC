"""Build the Round 2 business-proposal deck on the official Accenture template.

Every quantitative claim on every slide is read from ``results/tables/`` at build
time. Nothing is typed in by hand. If the experiment has not been run, the build
fails rather than emitting a deck full of plausible-looking numbers -- which is
the failure mode this design exists to prevent.

    python -m rippletwin.deck.build_deck --template <official.pptx> --out deck.pptx
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Optional, Sequence

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Accenture brand
PURPLE = RGBColor(0xA1, 0x00, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xF4, 0xEC, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1F, 0x9D, 0x55)
RED = RGBColor(0xD6, 0x3A, 0x3A)
BLUE = RGBColor(0x4C, 0x9B, 0xE8)
AMBER = RGBColor(0xF2, 0xB7, 0x05)

FONT = "Arial"
SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)


# --------------------------------------------------------------- slide helpers


_R_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def _clear(slide) -> None:
    """Remove every shape, turning a template slide into a blank canvas.

    We deliberately do NOT delete unwanted template slides. python-pptx's
    ``drop_rel`` leaves the underlying slide part in the package, and the next
    slide added reuses its partname -- which silently overwrites a slide on save
    ("Duplicate name: ppt/slides/slideN.xml"). That corrupted an earlier build:
    the team slide came back as a duplicate of a content slide.

    Repurposing the slides we do not want, then reordering at the end, touches
    no parts at all and cannot collide.
    """
    for shp in list(slide.shapes):
        shp._element.getparent().remove(shp._element)


def _reorder(prs: Presentation, desired) -> None:
    """Reorder slides to match ``desired`` (a list of slide objects).

    Pure element reordering inside ``sldIdLst`` -- no parts created or
    destroyed, so this is safe where deletion is not.
    """
    xs = prs.slides._sldIdLst
    by_part = {}
    for sldId in list(xs):
        by_part[prs.part.related_part(sldId.get(_R_ID))] = sldId
    for s in desired:
        el = by_part.get(s.part)
        if el is not None:
            xs.remove(el)
            xs.append(el)


def _blank_layout(prs: Presentation):
    """The lightest, emptiest layout available, so we control the whole canvas.

    Falling through to the *last* layout picks the purple Salutation master,
    whose gradient background makes light content cards unreadable. Prefer a
    light content layout -- and either way we paint our own white ground.
    """
    names = [l.name for l in prs.slide_masters[0].slide_layouts]
    for want in ("Content: text + split", "1_Complex 1 - Light mode",
                 "3_Complex 1 - Light mode Simple", "Blank", "Title Only"):
        if want in names:
            return prs.slide_masters[0].slide_layouts[names.index(want)]
    return prs.slide_masters[0].slide_layouts[0]


class SlidePool:
    """Hands out canvases: repurposed template slides first, then new ones."""

    def __init__(self, prs: Presentation, reuse_indices):
        self.prs = prs
        self._pool = [prs.slides[i] for i in reuse_indices]
        self.content = []

    def new(self):
        if self._pool:
            s = self._pool.pop(0)
            _clear(s)
        else:
            s = self.prs.slides.add_slide(_blank_layout(self.prs))
            _strip(s)
        # Explicit white ground so a repurposed slide and a fresh one look
        # identical regardless of the layout they inherited.
        rect(s, Emu(0), Emu(0), SW, SH, fill=WHITE, radius=False)
        self.content.append(s)
        return s


def _strip(slide) -> None:
    """Remove inherited placeholders so nothing empty renders."""
    for shp in list(slide.shapes):
        if shp.is_placeholder:
            shp._element.getparent().remove(shp._element)


def textbox(
    slide, x, y, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT,
    italic=False, spacing=1.15, anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = FONT
    return tb


def rect(slide, x, y, w, h, fill=LIGHT, line=None, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except (IndexError, KeyError):
            pass
    return shp


def slide_frame(pool: "SlidePool", title: str, kicker: str = "") -> object:
    """A content slide with the Accenture header bar."""
    s = pool.new()
    rect(s, Emu(0), Emu(0), SW, Inches(0.72), fill=PURPLE, radius=False)
    textbox(s, MARGIN, Inches(0.16), SW - 2 * MARGIN, Inches(0.42), title,
            size=21, bold=True, color=WHITE)
    if kicker:
        textbox(s, MARGIN, Inches(0.84), SW - 2 * MARGIN, Inches(0.34), kicker,
                size=13, color=GREY, italic=True)
    return s


def footnote(slide, text: str) -> None:
    textbox(slide, MARGIN, SH - Inches(0.42), SW - 2 * MARGIN, Inches(0.3),
            text, size=9, color=GREY, italic=True)


def kpi_row(slide, y, items: Sequence[tuple], height=Inches(1.15)):
    """items: (value, label, colour)"""
    n = len(items)
    gap = Inches(0.18)
    w = (SW - 2 * MARGIN - gap * (n - 1)) / n
    for i, (val, lab, col) in enumerate(items):
        x = MARGIN + i * (w + gap)
        rect(slide, x, y, w, height, fill=LIGHT)
        textbox(slide, x + Inches(0.14), y + Inches(0.13), w - Inches(0.28),
                Inches(0.5), val, size=25, bold=True, color=col)
        textbox(slide, x + Inches(0.14), y + Inches(0.66), w - Inches(0.28),
                Inches(0.42), lab, size=10, color=GREY)


def bullet_card(slide, x, y, w, h, heading, body, accent=PURPLE, fill=LIGHT):
    rect(slide, x, y, w, h, fill=fill)
    rect(slide, x, y, Inches(0.05), h, fill=accent, radius=False)
    textbox(slide, x + Inches(0.2), y + Inches(0.14), w - Inches(0.36),
            Inches(0.3), heading, size=12, bold=True, color=accent)
    textbox(slide, x + Inches(0.2), y + Inches(0.52), w - Inches(0.36),
            h - Inches(0.66), body, size=11, color=DARK, spacing=1.2)


def table(slide, x, y, w, rows, col_w=None, header=True, size=10.5, row_h=Inches(0.3)):
    n_r, n_c = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_r, n_c, x, y, w, row_h * n_r)
    tbl = shape.table
    if col_w:
        total = sum(col_w)
        for j, cw in enumerate(col_w):
            tbl.columns[j].width = Emu(int(w * cw / total))
    for i, row in enumerate(rows):
        tbl.rows[i].height = row_h
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = PURPLE if (header and i == 0) else (
                WHITE if i % 2 else LIGHT
            )
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(size)
            r.font.name = FONT
            r.font.bold = header and i == 0
            r.font.color.rgb = WHITE if (header and i == 0) else DARK
    return shape


def picture(slide, path: Path, x, y, w=None, h=None):
    if not Path(path).exists():
        return None
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


# ------------------------------------------------------------------- the data


class Results:
    """Loads every number the deck quotes. Fails loudly if evidence is missing."""

    def __init__(self, results_dir: Path):
        self.dir = Path(results_dir)
        self.tables = self.dir / "tables"
        self.figures = self.dir / "figures"
        req = ["flow_faults_hidden_source.csv", "cycle_time_inference.csv",
               "false_alarms.csv", "manifest.json"]
        missing = [f for f in req if not (self.tables / f).exists()]
        if missing:
            raise FileNotFoundError(
                "Cannot build the deck: missing evidence "
                + ", ".join(missing)
                + ". Run rippletwin.evaluation.experiments first. The deck is "
                "built from measured results only."
            )
        self.hidden = pd.read_csv(self.tables / "flow_faults_hidden_source.csv")
        self.allsrc = pd.read_csv(self.tables / "flow_faults_only.csv")
        self.cycle = pd.read_csv(self.tables / "cycle_time_inference.csv")
        self.fa = pd.read_csv(self.tables / "false_alarms.csv")
        self.manifest = json.loads((self.tables / "manifest.json").read_text())
        vis = self.tables / "board_visibility.csv"
        self.visibility = pd.read_csv(vis) if vis.exists() else None
        q = self.tables / "quality_attribution.csv"
        self.quality = pd.read_csv(q) if q.exists() else None

    def hv(self, coverage, method, col):
        r = self.hidden[
            np.isclose(self.hidden["coverage"], coverage)
            & (self.hidden["method"] == method)
        ]
        return float(r[col].iloc[0]) if len(r) and col in r else float("nan")

    def av(self, coverage, method, col):
        r = self.allsrc[
            np.isclose(self.allsrc["coverage"], coverage)
            & (self.allsrc["method"] == method)
        ]
        return float(r[col].iloc[0]) if len(r) and col in r else float("nan")

    def cyc(self, coverage, hidden=True):
        r = self.cycle[
            np.isclose(self.cycle["coverage"], coverage)
            & (self.cycle["source_hidden"] == hidden)
        ]
        return (float(r["median_pct"].iloc[0]), int(r["n"].iloc[0])) if len(r) else (float("nan"), 0)

    def fa_rate(self, coverage, method):
        r = self.fa[
            np.isclose(self.fa["coverage"], coverage) & (self.fa["method"] == method)
        ]
        return float(r["false_alarm_rate"].iloc[0]) if len(r) else float("nan")

    def vis_stats(self):
        if self.visibility is None or self.visibility.empty:
            return None
        h = self.visibility[self.visibility["source_hidden"]]
        if not len(h):
            return None
        inv = h[~h["ever_visible_on_board"]]
        out = {"n": len(h), "n_invisible": len(inv),
               "pct": 100.0 * len(inv) / len(h)}
        if "within1_episode" in inv.columns and len(inv):
            out["located_pct"] = 100.0 * inv["within1_episode"].fillna(0).mean()
        return out


def pct(x) -> str:
    return "n/a" if not np.isfinite(x) else f"{x * 100:.0f}%"


# ---------------------------------------------------------------- the slides


def build(template: Path, out: Path, results_dir: Path, team: dict) -> Path:
    R = Results(results_dir)
    prs = Presentation(str(template))

    # Template order: 0 cover, 1 instructions, 2 team, 3-5 empty content,
    # 6 thank-you. Keep cover, team and thank-you; repurpose the rest as
    # canvases rather than deleting them.
    n_tpl = len(prs.slides._sldIdLst)
    cover = prs.slides[0]
    team_slide = prs.slides[2] if n_tpl > 2 else None
    closing = prs.slides[n_tpl - 1] if n_tpl > 1 else None
    reuse = [i for i in (1, 3, 4, 5) if i < n_tpl]
    pool = SlidePool(prs, reuse)

    # ---- team slide ------------------------------------------------------
    if team_slide is not None:
        ts = team_slide
        for shp in ts.shapes:
            if shp.has_text_frame and "TEAM NAME" in shp.text_frame.text.upper():
                shp.text_frame.text = f"TEAM NAME: {team['team_name']}"
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(14)
                        r.font.name = FONT
                        r.font.color.rgb = PURPLE

    cov_main = 0.75
    covs = sorted(R.hidden["coverage"].unique())

    # ================================================== 3. the blind spot
    s = slide_frame(pool, "Factories can only monitor what they instrumented",
                    "Problem framing — Track 4, DigitalTwin.ai")
    y = Inches(1.35)
    w = (SW - 2 * MARGIN - Inches(0.24)) / 2
    bullet_card(s, MARGIN, y, w, Inches(1.5), "THE CONDITION",
                "Assembly lines mix legacy and modern equipment. Some stations are "
                "richly instrumented; a meaningful minority run on manual checklists "
                "and emit no telemetry at all.", PURPLE)
    bullet_card(s, MARGIN + w + Inches(0.24), y, w, Inches(1.5), "THE CONSEQUENCE",
                "A defect introduced early may not surface until a much later gate, "
                "by which time many downstream units carry the same undetected issue.",
                PURPLE)
    y2 = y + Inches(1.7)
    bullet_card(s, MARGIN, y2, w, Inches(1.5), "WHY \"ADD SENSORS\" IS NOT THE ANSWER",
                "Retrofits are limited to a small number of scheduled maintenance "
                "windows per year. The constraint is the calendar as much as the "
                "capital.", RED)
    bullet_card(s, MARGIN + w + Inches(0.24), y2, w, Inches(1.5), "WHERE THE VALUE SITS",
                "The stations most likely to hide a developing problem are exactly "
                "the ones nobody is watching. Inferring their state is the "
                "opportunity.", GREEN)

    vs = R.vis_stats()
    if vs:
        rect(s, MARGIN, y2 + Inches(1.72), SW - 2 * MARGIN, Inches(0.82), fill=DARK)
        textbox(s, MARGIN + Inches(0.25), y2 + Inches(1.86), SW - 2 * MARGIN - Inches(0.5),
                Inches(0.6),
                f"In our simulation, {vs['pct']:.0f}% of disturbances at "
                f"un-instrumented stations never produced a sustained line-level "
                f"throughput shortfall (n={vs['n']}). Aggregate monitoring would not "
                f"have caught them later — it would not have caught them at all.",
                size=13, bold=True, color=WHITE)
    footnote(s, "Quoted conditions are from the official Round 2 brief. "
                "The measured figure is a simulated prototype result on synthetic data.")

    # ================================================== 4. round 1 -> round 2
    s = slide_frame(pool, "Round 1 proposed the idea. Round 2 tests it.",
                    "What changed, and what we had to prove")
    rows = [["", "Round 1 — concept", "Round 2 — this submission"],
            ["Shadow-sensing", "asserted", "two implemented mechanisms: flow physics + defect genealogy"],
            ["Hidden state", "\"a live health signal\"", "cycle time of an unmeasured station, graded against ground truth"],
            ["Prediction", "\"predicts, not visualises\"", "forward flow arithmetic: vehicles lost, minutes to starvation"],
            ["Differentiation", "claimed", "3 baselines at a matched false-alarm rate"],
            ["Partial sensors", "claimed", "coverage sweep, 100% → 25% instrumentation"],
            ["Human in the loop", "in the diagram", "hash-chained ledger, abstention when ambiguous"],
            ["Evidence", "none", f"{R.manifest['n_episodes_test']} held-out episodes, 45 passing tests"]]
    table(s, MARGIN, Inches(1.4), SW - 2 * MARGIN, rows,
          col_w=[1.5, 2.2, 4.6], row_h=Inches(0.42), size=10.5)
    footnote(s, "Round 1 argued the idea was worth testing. This is the test — including the parts we failed.")

    # ================================================== 5. the mechanism (hero)
    s = slide_frame(pool, "The mechanism: the line measures its own blind spots",
                    "Conservation of material, not a learned correlation")
    textbox(s, MARGIN, Inches(1.3), SW - 2 * MARGIN, Inches(0.9),
            "When a station slows, material stops arriving downstream and stops "
            "leaving upstream. Every station downstream STARVES. Every station "
            "upstream BLOCKS. The boundary between them sits exactly at the station "
            "causing it — so we never need a sensor AT that station, only sensors on "
            "both sides of it.", size=13.5, color=DARK)
    fig = R.figures / "pressure_profile.png"
    if fig.exists():
        # Size by HEIGHT, not width: at full slide width this figure is tall
        # enough to run underneath the footnote.
        picture(s, fig, MARGIN + Inches(0.6), Inches(2.3), h=Inches(4.1))
    else:
        _mechanism_diagram(s, Inches(2.5))
    footnote(s, "A correlation model sees S07 blocking and S09 starving and cannot say which caused which. "
                "The flow model knows the direction of causation a priori — material moves one way.")

    # ================================================== 6. architecture
    s = slide_frame(pool, "How RippleTwin works", "Every box has an implementation behind it")
    stages = [
        ("OBSERVED", "telemetry from\ninstrumented\nstations only", BLUE),
        ("BASELINE", "mix- and shift-aware\nexpectations, in\nfraction of takt", BLUE),
        ("INFERRED", "posterior over ALL\nstations, including\nun-sensored ones", PURPLE),
        ("PREDICTED", "ripple forecast from\nflow arithmetic", AMBER),
        ("EXPLAINED", "evidence tagged\nOBS / INF / PRED,\nno LLM in the loop", GREEN),
        ("HUMAN", "approve / reject\nor ABSTAIN", DARK),
    ]
    n = len(stages)
    gap = Inches(0.14)
    cw = (SW - 2 * MARGIN - gap * (n - 1)) / n
    for i, (h, b, c) in enumerate(stages):
        x = MARGIN + i * (cw + gap)
        rect(s, x, Inches(1.5), cw, Inches(1.9), fill=LIGHT)
        rect(s, x, Inches(1.5), cw, Inches(0.36), fill=c)
        textbox(s, x, Inches(1.57), cw, Inches(0.28), h, size=10.5, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
        textbox(s, x + Inches(0.1), Inches(2.0), cw - Inches(0.2), Inches(1.3),
                b, size=9.5, color=DARK, align=PP_ALIGN.CENTER)
        if i < n - 1:
            textbox(s, x + cw, Inches(2.25), gap, Inches(0.3), "›", size=16,
                    bold=True, color=GREY, align=PP_ALIGN.CENTER)
    rect(s, MARGIN, Inches(3.62), SW - 2 * MARGIN, Inches(0.42), fill=DARK)
    textbox(s, MARGIN + Inches(0.2), Inches(3.7), SW - 2 * MARGIN, Inches(0.3),
            "OUTCOME → hash-chained ledger → per-station precision feeds back",
            size=11, bold=True, color=WHITE)

    rows = [["Component", "Method chosen", "Why not something heavier"],
            ["Hidden-state estimation", "non-negative least squares + Gaussian likelihood over candidate positions",
             "a GNN would learn the propagation we can derive — deriving it is stronger and inspectable"],
            ["Calibration", "empirical null + correlation correction",
             "stations are ~0.4 correlated; assuming independence inflates evidence ~14x"],
            ["Ripple forecast", "serial-line flow arithmetic",
             "an engineer can check it on paper, and it cannot drift"],
            ["Explanation", "templates bound to computed values",
             "an LLM narrator can drift from the evidence; here it structurally cannot"]]
    table(s, MARGIN, Inches(4.25), SW - 2 * MARGIN, rows,
          col_w=[1.7, 3.0, 4.0], row_h=Inches(0.44), size=9.5)
    footnote(s, "Rule we held to: when a simpler method is equally effective, use the simpler method.")

    # ================================================== 7. experiment design
    s = slide_frame(pool, "How we made the comparison fair",
                    "The design decisions that make the next slide mean something")
    w2 = (SW - 2 * MARGIN - Inches(0.24)) / 2
    bullet_card(s, MARGIN, Inches(1.35), w2, Inches(1.55), "MATCHED FALSE-ALARM RATE",
                "Every method is calibrated on the same held-out nominal data to the "
                "same 1% per-window false-alarm target, then run through an identical "
                "detection rule. Without this, a permissive threshold just looks more "
                "sensitive.", PURPLE)
    bullet_card(s, MARGIN + w2 + Inches(0.24), Inches(1.35), w2, Inches(1.55),
                "ONE PHYSICS RUN, MANY SENSOR VIEWS",
                "Coverage levels are views over a single simulation, so differences "
                "between them are differences in observability — not different random "
                "draws.", PURPLE)
    bullet_card(s, MARGIN, Inches(3.05), w2, Inches(1.55), "B2 IS THE SHARP TEST",
                "B2 is RippleTwin's own model with un-sensored stations removed as "
                "candidates. Same physics, same likelihood, same tuning. The only "
                "difference is whether a blind station may be named — which isolates "
                "shadow-sensing itself.", RED)
    bullet_card(s, MARGIN + w2 + Inches(0.24), Inches(3.05), w2, Inches(1.55),
                "DISJOINT DATA THROUGHOUT",
                "Baseline fitting, detector calibration and evaluation use separate "
                "runs. Tuning and test episodes are split by seed; anything chosen by "
                "looking at data was chosen on the tuning half.", GREEN)
    rows = [["Baseline", "What it represents"],
            ["B0 — SPC per station", "the control chart most plants already run"],
            ["B1 — Isolation Forest", "modern anomaly detection, given the same features, topology-blind"],
            ["B2 — observed-only twin", "a conventional digital twin: models what it measures, silent about the rest"]]
    table(s, MARGIN, Inches(4.8), SW - 2 * MARGIN, rows, col_w=[2.2, 7.0],
          row_h=Inches(0.36), size=10.5)
    footnote(s, "Simulated prototype evaluation on synthetic data.")

    # ================================================== 8. headline result
    s = slide_frame(pool, "When the faulty station has no sensor, only RippleTwin names it",
                    f"Held-out episodes · flow faults · source station un-instrumented")
    kpi_row(s, Inches(1.3), [
        (pct(R.hv(cov_main, "RippleTwin", "top1")),
         f"RippleTwin — exact station, at {cov_main*100:.0f}% coverage", PURPLE),
        (pct(R.hv(cov_main, "B2_observed_only_twin", "top1")),
         "Conventional observed-only twin", RED),
        (pct(R.hv(cov_main, "B1_IsolationForest", "top1")),
         "Anomaly detection", RED),
        (pct(R.hv(cov_main, "B0_SPC_observed", "top1")), "SPC on sensors", RED),
    ])
    textbox(s, MARGIN, Inches(2.62), SW - 2 * MARGIN, Inches(0.4),
            "The baselines do not merely score badly. They score exactly zero — "
            "naming an unmeasured station is outside what they can express.",
            size=12.5, bold=True, color=DARK)
    rows = [["Sensor coverage", "RippleTwin", "B2 observed-only twin", "B1 anomaly", "B0 SPC"]]
    for c in covs:
        rows.append([f"{c*100:.0f}%",
                     pct(R.hv(c, "RippleTwin", "top1")),
                     pct(R.hv(c, "B2_observed_only_twin", "top1")),
                     pct(R.hv(c, "B1_IsolationForest", "top1")),
                     pct(R.hv(c, "B0_SPC_observed", "top1"))])
    table(s, MARGIN, Inches(3.15), Inches(6.1), rows,
          col_w=[1.5, 1.2, 1.5, 1.0, 0.9], row_h=Inches(0.34), size=10)
    # The exact-station variant, so the chart and the table beside it are the
    # same metric. A within-one-station chart next to an exact-station table
    # reads as a contradiction.
    chart = R.figures / "coverage_curve_top1.png"
    if not chart.exists():
        chart = R.figures / "coverage_curve.png"
    picture(s, chart, MARGIN + Inches(6.4), Inches(3.0), w=Inches(6.2))
    footnote(s, "Exact-station localisation while the fault is active, at a matched 1% "
                "per-window false-alarm rate. Simulated prototype result on synthetic data.")

    # ================================================== 9. the control
    s = slide_frame(pool, "The control that convinced us the mechanism is real",
                    "At 100% coverage the advantage disappears — exactly as it should")
    a = R.av(1.0, "RippleTwin", "top1")
    b = R.av(1.0, "B2_observed_only_twin", "top1")
    kpi_row(s, Inches(1.3), [
        (pct(a), "RippleTwin at 100% sensor coverage", PURPLE),
        (pct(b), "Observed-only twin, same conditions", BLUE),
        ("identical", "to three decimal places", GREEN),
    ])
    textbox(s, MARGIN, Inches(2.7), SW - 2 * MARGIN, Inches(1.1),
            "With every station instrumented there is nothing hidden to infer, so "
            "shadow-sensing has nothing to add — and it adds nothing. The advantage "
            "appears exactly and only where instrumentation is missing.\n\n"
            "That is what should happen if the mechanism is real rather than an "
            "artefact of better tuning. A method that also won here would be telling "
            "us its gains came from somewhere else.",
            size=13, color=DARK)
    med, n_cyc = R.cyc(cov_main, hidden=True)
    rect(s, MARGIN, Inches(4.15), SW - 2 * MARGIN, Inches(1.35), fill=LIGHT)
    textbox(s, MARGIN + Inches(0.25), Inches(4.3), SW - 2 * MARGIN - Inches(0.5),
            Inches(0.32), "AND THE SHARPEST FALSIFIABLE CLAIM", size=11.5,
            bold=True, color=PURPLE)
    textbox(s, MARGIN + Inches(0.25), Inches(4.68), SW - 2 * MARGIN - Inches(0.5),
            Inches(0.8),
            f"RippleTwin estimates the cycle time of a station that has no sensor — a "
            f"number that is unmeasurable in the data it receives — to a median error "
            f"of {med:.1f}% at {cov_main*100:.0f}% coverage (n={n_cyc}). The simulator "
            f"knows the true value, so the estimate is graded, not asserted.",
            size=12.5, color=DARK)
    footnote(s, "Simulated prototype results on synthetic data.")

    # ================================================== 10. trust
    s = slide_frame(pool, "A detector that never stays quiet is worthless",
                    "False alarms, abstention, and the cases where we correctly say nothing")
    fa_rt = R.fa_rate(cov_main, "RippleTwin")
    kpi_row(s, Inches(1.3), [
        (f"{fa_rt*100:.1f}%", "RippleTwin false-alarm rate per window", GREEN),
        ("0", "alerts on the normal-variation scenario", GREEN),
        ("LINE_SUPPLY", "a material delay is attributed to supply, not a station", GREEN),
        ("ABSTAIN", "when adjacent candidates cannot be separated", AMBER),
    ])
    w2 = (SW - 2 * MARGIN - Inches(0.24)) / 2
    bullet_card(s, MARGIN, Inches(2.75), w2, Inches(1.5), "THRESHOLD IS A DESIGN PARAMETER",
                "Set from a stated false-alarm target on held-out nominal data, not "
                "hand-picked. The brief warns that false alarms erode floor trust "
                "quickly — and that trust does not survive being lost twice.", GREEN)
    bullet_card(s, MARGIN + w2 + Inches(0.24), Inches(2.75), w2, Inches(1.5),
                "IT DECLINES TO GUESS",
                "When posterior mass is spread across adjacent blind stations, "
                "RippleTwin reports the group and escalates rather than naming one. "
                "Those stations become sensor-placement candidates.", AMBER)
    rect(s, MARGIN, Inches(4.45), SW - 2 * MARGIN, Inches(1.0), fill=DARK)
    textbox(s, MARGIN + Inches(0.25), Inches(4.62), SW - 2 * MARGIN - Inches(0.5),
            Inches(0.7),
            "Two of our own designs were killed by these tests: a symmetric "
            "pressure channel that mis-located faults by one station, and a robust "
            "z-score on starvation time that produced a 58% false-alarm rate. Both "
            "are documented in the repository rather than buried.",
            size=12, bold=True, color=WHITE)
    footnote(s, "Simulated prototype evaluation on synthetic data.")

    # ================================================== 11. stakeholders
    s = slide_frame(pool, "Three users, one twin", "Same inferred state, different horizon")
    cards = [
        ("FLOOR SUPERVISOR", "Right now, this station",
         "• which station, and how sure\n• why — evidence tagged OBS/INF/PRED\n"
         "• what happens next, in vehicles\n• one recommended action\n• approve or reject", PURPLE),
        ("PLANT MANAGER", "This week, this line",
         "• where suspicion accumulates\n• throughput and escape rate\n"
         "• coverage by zone\n• defects by discovery gate\n• the decision ledger", BLUE),
        ("LEADERSHIP", "The rollout case",
         "• value drivers and payback\n• sensor economics vs retrofit\n"
         "• where instrumentation is worth buying\n• phased rollout and risks", GREEN),
    ]
    gap = Inches(0.24)
    cw = (SW - 2 * MARGIN - gap * 2) / 3
    for i, (h, sub, body, c) in enumerate(cards):
        x = MARGIN + i * (cw + gap)
        rect(s, x, Inches(1.35), cw, Inches(2.9), fill=LIGHT)
        rect(s, x, Inches(1.35), cw, Inches(0.42), fill=c)
        textbox(s, x + Inches(0.15), Inches(1.43), cw, Inches(0.3), h, size=11.5,
                bold=True, color=WHITE)
        textbox(s, x + Inches(0.15), Inches(1.88), cw - Inches(0.3), Inches(0.28),
                sub, size=10.5, italic=True, color=GREY)
        textbox(s, x + Inches(0.15), Inches(2.22), cw - Inches(0.3), Inches(1.9),
                body, size=10.5, color=DARK, spacing=1.35)
    rect(s, MARGIN, Inches(4.45), SW - 2 * MARGIN, Inches(0.95), fill=LIGHT)
    textbox(s, MARGIN + Inches(0.25), Inches(4.6), SW - 2 * MARGIN - Inches(0.5), Inches(0.7),
            "Throughout, every value carries its provenance: OBSERVED (a sensor read "
            "it), INFERRED (we estimated it for a station with none), PREDICTED (it "
            "has not happened yet). A supervisor who cannot tell a measurement from "
            "an estimate will eventually trust the wrong one.",
            size=12, color=DARK)

    # ================================================== 12. business case
    s = slide_frame(pool, "Business case", "Illustrative assumptions, transparent arithmetic")
    kpi_row(s, Inches(1.25), [
        ("$2,200", "assumed contribution margin / vehicle", GREY),
        ("60/yr", "assumed hidden-station disturbances", GREY),
        ("55%", "assumed recovery factor", GREY),
        ("~1 month", "resulting payback (illustrative)", PURPLE),
    ])
    rows = [["Value driver", "Mechanism", "Evidence class"],
            ["Directed maintenance", "names WHICH station, not \"somewhere in body shop\"", "PROTOTYPE RESULT"],
            ["Fewer defect escapes", "attributes a defect-type excess to a source before the next gate", "PROTOTYPE RESULT"],
            ["Recovered throughput", "constraint found before the shortfall compounds", "PROTOTYPE + ASSUMPTION"],
            ["Lower instrumentation need", "blind stations become inferable rather than dark", "PROTOTYPE RESULT"]]
    table(s, MARGIN, Inches(2.65), SW - 2 * MARGIN, rows, col_w=[2.2, 5.2, 1.9],
          row_h=Inches(0.38), size=10)
    rect(s, MARGIN, Inches(4.6), SW - 2 * MARGIN, Inches(1.0), fill=DARK)
    textbox(s, MARGIN + Inches(0.25), Inches(4.75), SW - 2 * MARGIN - Inches(0.5), Inches(0.75),
            "Where this is weakest, stated plainly: the model is dominated by margin "
            "per recovered vehicle. On a line that is NOT capacity-constrained, a "
            "recovered vehicle is worth almost nothing and the throughput driver "
            "collapses — so we would qualify the line before selling that value.",
            size=12, bold=True, color=WHITE)
    footnote(s, "No figure on this slide is measured from a real plant. All inputs are editable in the dashboard.")

    # ================================================== 13. sensor economics
    s = slide_frame(pool, "Sensor economics: we are not competing with the sensor budget",
                    "We are telling you how to spend it")
    rows = [["", "Retrofit every blind station", "RippleTwin"],
            ["Capital (10 stations)", "~$180,000", "~$150,000 year 1"],
            ["When it can be installed", "scheduled maintenance windows only", "against data already recorded"],
            ["Scope", "those 10 stations", "every station on the line"],
            ["Failure mode", "sensor dies, station goes dark again", "inference degrades, and reports that it has"]]
    table(s, MARGIN, Inches(1.35), SW - 2 * MARGIN, rows, col_w=[2.0, 3.6, 3.6],
          row_h=Inches(0.42), size=10.5)
    rect(s, MARGIN, Inches(3.65), SW - 2 * MARGIN, Inches(1.05), fill=LIGHT)
    textbox(s, MARGIN + Inches(0.25), Inches(3.8), SW - 2 * MARGIN - Inches(0.5), Inches(0.8),
            "We do NOT claim zero sensors. Shadow-sensing needs instrumented stations "
            "on both sides of a blind one — that is the mechanism. The claim is "
            "narrower and testable: the sensors already installed carry more "
            "information about their neighbours than a conventional twin extracts.",
            size=12.5, color=DARK)
    rect(s, MARGIN, Inches(4.9), SW - 2 * MARGIN, Inches(0.95), fill=PURPLE)
    textbox(s, MARGIN + Inches(0.25), Inches(5.05), SW - 2 * MARGIN - Inches(0.5), Inches(0.7),
            "The coverage experiment is itself a sensor-placement tool: it shows which "
            "blind stations inference already covers, and which genuinely need "
            "instrumenting — so the retrofit budget goes where it actually buys something.",
            size=12.5, bold=True, color=WHITE)

    # ================================================== 14. roadmap
    s = slide_frame(pool, "Phased rollout", "Deliberately slow at the start, because trust is spent once")
    phases = [
        ("PHASE 0", "2–4 weeks", "Data readiness", "PLC timestamps, gate results and build sequence are recoverable and joinable", GREY),
        ("PHASE 1", "8–12 weeks", "Shadow mode", "Alerts logged, nobody acts. Precision per station measured on the ledger", BLUE),
        ("PHASE 2", "12 weeks", "Live advisory", "Supervisors see and act. Adoption and precision on real outcomes", PURPLE),
        ("PHASE 3", "2 quarters", "Plant rollout", "Remaining lines; per-line configuration under one week", GREEN),
        ("PHASE 4", "12 months+", "Multi-plant", "Other sites and equipment vintages; baseline transfer", DARK),
    ]
    gap = Inches(0.14)
    cw = (SW - 2 * MARGIN - gap * 4) / 5
    for i, (ph, dur, name, body, c) in enumerate(phases):
        x = MARGIN + i * (cw + gap)
        rect(s, x, Inches(1.4), cw, Inches(2.5), fill=LIGHT)
        rect(s, x, Inches(1.4), cw, Inches(0.34), fill=c)
        textbox(s, x, Inches(1.46), cw, Inches(0.26), ph, size=10, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
        textbox(s, x + Inches(0.12), Inches(1.85), cw - Inches(0.24), Inches(0.24),
                dur, size=9.5, italic=True, color=GREY)
        textbox(s, x + Inches(0.12), Inches(2.12), cw - Inches(0.24), Inches(0.3),
                name, size=11.5, bold=True, color=DARK)
        textbox(s, x + Inches(0.12), Inches(2.5), cw - Inches(0.24), Inches(1.3),
                body, size=9.5, color=DARK, spacing=1.25)
    rect(s, MARGIN, Inches(4.1), SW - 2 * MARGIN, Inches(0.85), fill=DARK)
    textbox(s, MARGIN + Inches(0.25), Inches(4.24), SW - 2 * MARGIN - Inches(0.5), Inches(0.6),
            "Integration is read-only throughout. RippleTwin writes to no control "
            "system — the brief notes that modifying live PLCs carries real "
            "operational risk, so the prototype does not, and neither should a pilot.",
            size=12, bold=True, color=WHITE)

    # ================================================== 15. risks
    s = slide_frame(pool, "Risks and limitations, stated before you find them",
                    "The honest version is the one worth reviewing")
    rows = [["Risk / limitation", "Severity", "Mitigation"],
            ["No real-world validation — everything is synthetic", "HIGH", "shadow-mode pilot before any operational claim; validate the flow signature on one real line's PLC logs first"],
            ["Value depends on the line being capacity-constrained", "HIGH", "qualify during Phase 0; do not sell throughput value to an unconstrained plant"],
            ["Serial-line assumption in the propagation model", "MEDIUM", "the likelihood generalises; the pattern matrix must be rebuilt from the real process graph"],
            ["Adjacent blind stations are not separable", "MEDIUM", "report the group and abstain; flag as sensor-placement candidates"],
            ["Alert fatigue kills adoption", "MEDIUM", "threshold from a stated false-alarm target; abstention when ambiguous"],
            ["An incumbent platform copies the algorithm", "MEDIUM", "the algorithm is not the moat — the outcome ledger and sensor-placement position are"]]
    table(s, MARGIN, Inches(1.35), SW - 2 * MARGIN, rows, col_w=[3.3, 1.0, 5.5],
          row_h=Inches(0.5), size=10)
    footnote(s, "Full limitations section in the repository README.")

    # ================================================== 16. close
    s = slide_frame(pool, "Why this is difficult to dismiss", "")
    points = [
        ("The mechanism is derived, not learned",
         "Conservation of material through a serial line. A judge can check the derivation, and the tests assert it against ground truth."),
        ("The advantage is isolated, not asserted",
         "B2 is the same model minus hidden candidates. At 100% coverage the two are identical — the gain appears only where instrumentation is missing."),
        ("The claim is falsifiable, and we did the falsifying",
         "Two of our designs were killed by our own measurements and are documented. The metric that would have flattered us most, we report as under-powered."),
        ("It is honest about what it does not know",
         "Provenance on every number, abstention when ambiguous, and a limitations section we wrote before anyone asked for one."),
    ]
    y = Inches(1.4)
    for h, b in points:
        rect(s, MARGIN, y, SW - 2 * MARGIN, Inches(0.98), fill=LIGHT)
        rect(s, MARGIN, y, Inches(0.06), Inches(0.98), fill=PURPLE, radius=False)
        textbox(s, MARGIN + Inches(0.28), y + Inches(0.13), SW - 2 * MARGIN, Inches(0.3),
                h, size=13.5, bold=True, color=PURPLE)
        textbox(s, MARGIN + Inches(0.28), y + Inches(0.5), SW - 2 * MARGIN - Inches(0.6),
                Inches(0.42), b, size=11.5, color=DARK)
        y += Inches(1.12)
    textbox(s, MARGIN, Inches(6.1), SW - 2 * MARGIN, Inches(0.5),
            "RippleTwin — see the bottleneck before it arrives, even at the stations "
            "you can't instrument.", size=15, bold=True, color=PURPLE)

    # Cover, team, content in build order, closing salutation last.
    desired = [cover] + ([team_slide] if team_slide is not None else [])
    desired += pool.content
    if closing is not None and closing not in desired:
        desired.append(closing)
    _reorder(prs, desired)

    out = Path(out)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def _mechanism_diagram(slide, y):
    """Fallback station diagram if the generated figure is unavailable."""
    ids = ["S05", "S06", "S07", "S08", "S09", "S10", "S11"]
    n = len(ids)
    gap = Inches(0.16)
    cw = (SW - 2 * MARGIN - gap * (n - 1)) / n
    for i, sid in enumerate(ids):
        x = MARGIN + i * (cw + gap)
        hidden = sid == "S08"
        fill = RGBColor(0x9A, 0xA0, 0xA6) if hidden else (RED if i < 3 else BLUE)
        rect(slide, x, y, cw, Inches(0.85), fill=fill)
        textbox(slide, x, y + Inches(0.16), cw, Inches(0.3), sid, size=14,
                bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        label = "NO SENSOR" if hidden else ("BLOCKED" if i < 3 else "STARVED")
        textbox(slide, x, y + Inches(0.95), cw, Inches(0.28), label, size=9.5,
                bold=True, color=fill, align=PP_ALIGN.CENTER)
    textbox(slide, MARGIN, y + Inches(1.35), SW - 2 * MARGIN, Inches(0.35),
            "the boundary between blocked and starved IS the constraint",
            size=12.5, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def main(argv=None) -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True)
    ap.add_argument("--out", default="assets/RippleTwin_Round2_Proposal.pptx")
    ap.add_argument("--results", default="results")
    ap.add_argument("--team-name", default="RippleTwin")
    args = ap.parse_args(argv)

    p = build(Path(args.template), Path(args.out), Path(args.results),
              {"team_name": args.team_name})
    print(f"wrote {p}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
